"""Registre des extracteurs de texte par extension de fichier.

Le reste de l'outil ne connaît pas les types de fichiers spécifiquement : il
délègue l'extraction du texte à l'extracteur enregistré pour l'extension du
fichier. Ajouter un nouveau type de fichier ne demande qu'un nouvel
extracteur ici.

Les formats "texte brut" et ceux qui ne demandent que la bibliothèque
standard (html, xml, rtf, eml...) sont toujours disponibles. Les formats qui
nécessitent un paquet tiers (docx, xlsx/pptx, msg) l'importent en différé et
n'échouent qu'à l'usage, avec un message clair, pour ne jamais imposer une
dépendance lourde à qui ne s'en sert pas.
"""
from __future__ import annotations

import os
import re


def _extract_pdf(path: str) -> str:
    import pymupdf

    with pymupdf.open(path) as doc:
        text = " ".join(page.get_text() for page in doc)
    # Un PDF scanné (image du document, sans couche texte) donne une chaîne
    # vide ou uniquement des espaces ici — sans quoi il resterait
    # "illisible" (voir discover._build_bundle) alors que l'OCR (si activé,
    # config `ocr_enabled`) peut justement en tirer du texte exploitable.
    # Tenté seulement en dernier recours : l'extraction normale est
    # nettement plus rapide et déjà fiable pour un PDF avec texte natif.
    if not text.strip() and _ocr_config().ocr_enabled:
        return _extract_pdf_ocr(path)
    return text


def _extract_plain_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _strip_markup(markup_text: str) -> str:
    """Retire les balises HTML/XML d'un texte, ne garde que le contenu."""
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.chunks: list[str] = []

        def handle_data(self, data: str) -> None:
            self.chunks.append(data)

    parser = _TextExtractor()
    parser.feed(markup_text)
    return " ".join(parser.chunks)


def _extract_html(path: str) -> str:
    return _strip_markup(_extract_plain_text(path))


def _extract_xml(path: str) -> str:
    import xml.etree.ElementTree as ET

    try:
        tree = ET.parse(path)
        return " ".join(t for t in tree.getroot().itertext() if t)
    except ET.ParseError:
        # XML mal formé (ou en fait du HTML) : on retombe sur un nettoyage générique.
        return _strip_markup(_extract_plain_text(path))


_RTF_CONTROL_WORD = re.compile(r"\\[a-zA-Z]+-?\d* ?")
_RTF_HEX_ESCAPE = re.compile(r"\\'[0-9a-fA-F]{2}")
_RTF_BRACES = re.compile(r"[{}]")


def _extract_rtf(path: str) -> str:
    """Suppression grossière des groupes de contrôle RTF. Ne vise pas une
    fidélité parfaite, seulement à récupérer l'essentiel du texte."""
    raw = _extract_plain_text(path)
    text = _RTF_HEX_ESCAPE.sub(" ", raw)
    text = _RTF_CONTROL_WORD.sub(" ", text)
    text = _RTF_BRACES.sub(" ", text)
    return text


def _extract_eml(path: str) -> str:
    from email import policy
    from email.parser import BytesParser

    with open(path, "rb") as f:
        message = BytesParser(policy=policy.default).parse(f)

    parts = [message.get("subject", "") or ""]
    body = message.get_body(preferencelist=("plain", "html"))
    if body is not None:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            content = _strip_markup(content)
        parts.append(content)
    return "\n".join(parts)


def _extract_docx(path: str) -> str:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError(
            "La lecture des fichiers .docx nécessite le paquet optionnel python-docx : "
            "pip install -r requirements-docx.txt"
        ) from exc
    document = docx.Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError(
            "La lecture des fichiers .xlsx nécessite le paquet optionnel openpyxl : "
            "pip install -r requirements-office.txt"
        ) from exc
    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    values = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            values.extend(str(cell) for cell in row if cell is not None)
    return " ".join(values)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "La lecture des fichiers .pptx nécessite le paquet optionnel python-pptx : "
            "pip install -r requirements-office.txt"
        ) from exc
    presentation = Presentation(path)
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_odf(path: str) -> str:
    try:
        from odf import teletype
        from odf.opendocument import load
    except ImportError as exc:
        raise RuntimeError(
            "La lecture des fichiers OpenDocument (.odt/.ods/.odp) nécessite le paquet "
            "optionnel odfpy : pip install -r requirements-office.txt"
        ) from exc
    # `.body` est le point d'entrée générique du contenu quel que soit le
    # type de document (texte, tableur, présentation...) — plus fiable que de
    # distinguer chaque type par sa propre structure interne.
    document = load(path)
    return teletype.extractText(document.body)


def _ocr_config():
    from .config import get_config

    return get_config()


def _ocr_image(image) -> str:
    """OCR d'une image déjà chargée (objet PIL) via Tesseract — utilisé aussi
    bien pour un fichier image direct que pour une page de PDF scannée
    rasterisée (voir `_extract_pdf`). Le paquet Python `pytesseract` est
    installé par défaut (`requirements.txt`) ; seul le moteur Tesseract
    lui-même (binaire externe, pas un paquet pip) reste à installer
    séparément sur la machine — l'`ImportError` reste gérée ici en filet de
    sécurité (ex. installation manuelle sans ce paquet), mais ne devrait
    normalement jamais se produire. Lève `RuntimeError` avec un message
    clair si Tesseract est absent, plutôt que de laisser remonter une
    `TesseractNotFoundError` peu compréhensible."""
    try:
        import pytesseract
    except ImportError as exc:
        raise RuntimeError(
            "L'OCR nécessite le paquet pytesseract (normalement déjà installé par "
            "requirements.txt) : pip install pytesseract"
        ) from exc
    config = _ocr_config()
    if config.tesseract_cmd_path:
        pytesseract.pytesseract.tesseract_cmd = config.tesseract_cmd_path
    try:
        return pytesseract.image_to_string(image, lang="fra+eng")
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError(
            "L'OCR nécessite aussi le moteur Tesseract installé sur la machine (pas un "
            "paquet Python, un exécutable à part) — voir la section \"Formats de fichiers "
            "pris en charge\" du README pour les liens d'installation, ou précisez son "
            "chemin dans l'onglet Paramètres."
        ) from exc


def _extract_image(path: str) -> str:
    if not _ocr_config().ocr_enabled:
        raise RuntimeError(
            "La lecture des images nécessite l'OCR — activez-le dans l'onglet Paramètres "
            "(voir la config `ocr_enabled`)."
        )
    from PIL import Image

    with Image.open(path) as image:
        return _ocr_image(image)


def _extract_pdf_ocr(path: str) -> str:
    """Repli OCR pour un PDF scanné (sans couche texte) : rasterise chaque
    page (pymupdf, déjà une dépendance de base) puis reconnaît le texte page
    par page. Nettement plus lent que `_extract_pdf` — seulement tenté quand
    l'extraction normale n'a rien donné (voir `_extract_pdf`) et que l'OCR
    est activé (config `ocr_enabled`)."""
    import pymupdf
    from PIL import Image

    parts = []
    with pymupdf.open(path) as doc:
        for page in doc:
            # 200 DPI : compromis lisibilité/vitesse pour un scan de document
            # bureautique typique (la résolution par défaut, 72 DPI, est trop
            # basse pour une reconnaissance fiable des petits caractères).
            pixmap = page.get_pixmap(dpi=200)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            parts.append(_ocr_image(image))
    return "\n".join(parts)


def _extract_msg(path: str) -> str:
    try:
        import extract_msg
    except ImportError as exc:
        raise RuntimeError(
            "La lecture des fichiers .msg (Outlook) nécessite le paquet optionnel extract-msg : "
            "pip install -r requirements-msg.txt"
        ) from exc
    # `Message` garde le fichier .msg ouvert (accès OLE via olefile) tant
    # qu'il n'est pas fermé explicitement : sans le context manager, le
    # handle restait ouvert jusqu'au passage du ramasse-miettes, ce qui
    # pouvait faire échouer une opération ultérieure sur ce même fichier sous
    # Windows (copie dans dataset/, export vers le dossier de sortie...) —
    # le fichier semblait alors "verrouillé" par l'outil lui-même.
    with extract_msg.Message(path) as message:
        # `body` se replie déjà tout seul sur le corps RTF/HTML dé-encapsulé
        # si le message n'a pas de corps texte brut (voir extract_msg) ;
        # `sender` complète le sujet et le corps pour l'aperçu et la
        # catégorisation, comme le ferait un lecteur d'e-mails.
        #
        # Chaque champ est lu isolément : la désencapsulation RTF (utilisée
        # par `body` quand il n'y a pas de flux texte brut) est un point de
        # fragilité connu de la bibliothèque sur certains messages malformés
        # — un champ qui échoue ne doit pas priver l'extraction des autres,
        # encore exploitables, et faire perdre TOUT le texte du message.
        parts = []
        for attr in ("subject", "sender", "body"):
            try:
                value = getattr(message, attr)
            except Exception:
                value = None
            if value:
                parts.append(value)
        return "\n".join(parts)


# Extension (minuscule, avec le point) -> fonction d'extraction(path) -> texte.
EXTRACTORS = {
    # Documents
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".rtf": _extract_rtf,
    # Texte brut et données structurées lisibles telles quelles
    ".txt": _extract_plain_text,
    ".md": _extract_plain_text,
    ".csv": _extract_plain_text,
    ".tsv": _extract_plain_text,
    ".log": _extract_plain_text,
    ".json": _extract_plain_text,
    ".yaml": _extract_plain_text,
    ".yml": _extract_plain_text,
    ".ini": _extract_plain_text,
    ".cfg": _extract_plain_text,
    ".toml": _extract_plain_text,
    # Balisage
    ".html": _extract_html,
    ".htm": _extract_html,
    ".xml": _extract_xml,
    # Email
    ".eml": _extract_eml,
    ".msg": _extract_msg,
    # Bureautique (tableurs, présentations)
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".pptx": _extract_pptx,
    # OpenDocument (LibreOffice/OpenOffice)
    ".odt": _extract_odf,
    ".ods": _extract_odf,
    ".odp": _extract_odf,
    # Images (OCR requis, voir config `ocr_enabled` — non activé par défaut,
    # échoue sinon avec un message clair plutôt que de laisser un fichier
    # image se retrouver silencieusement sans texte)
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".tiff": _extract_image,
    ".bmp": _extract_image,
}

SUPPORTED_EXTENSIONS = tuple(sorted(EXTRACTORS))

# Sous-ensemble d'extensions image de EXTRACTORS ci-dessus — réutilisé par
# `discover._build_bundle` pour le moteur "image" (analyse visuelle par
# CLIP), qui ne doit considérer QUE des fichiers image, pas n'importe quel
# format pris en charge en général.
IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tiff", ".bmp"})


def is_supported(path: str) -> bool:
    return os.path.splitext(path)[1].lower() in EXTRACTORS


def extract_text_for(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(f"Type de fichier non pris en charge : {ext or path}")
    return extractor(path)
